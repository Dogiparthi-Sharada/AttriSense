# ---------------------------------------------------------------------------
# AttriSense — scripts/append_pitch_gallery.py
# ---------------------------------------------------------------------------
# Author : Sharada Dogiparthi <dogiparthi.sharada@gmail.com>
# Version: 1.0.0
# Date   : 2026-05-07
# License: MIT — see LICENSE in repo root.
# Copyright (c) 2026 Sharada Dogiparthi. All rights reserved.
# ---------------------------------------------------------------------------
"""Append a Pixel-pastel dashboard-gallery appendix to AttriSense_VP_Pitch.pptx.

Adds 8 appendix slides, each embedding one freshly-recaptured dashboard
screenshot (post dark-axis fix), so reviewers see the actual UI behind
the pitch rather than only bullet text.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

REPO = Path(__file__).resolve().parents[1]
PPTX = REPO / 'outputs' / 'AttriSense_VP_Pitch.pptx'
IMG = REPO / 'docs' / 'images'

# Pixel-pastel palette
CREAM = RGBColor(0xFB, 0xF7, 0xF2)
INK = RGBColor(0x3D, 0x3A, 0x35)
LAV = RGBColor(0xB8, 0xB5, 0xE1)
PEACH = RGBColor(0xF5, 0xC4, 0xA0)
SAGE = RGBColor(0xB5, 0xD4, 0xA8)
SAND = RGBColor(0xF5, 0xEF, 0xE6)

# (slug, headline, caption)
GALLERY = [
    ('executive',
     'Executive Overview',
     'Five surfaces, seven tabs, one dashboard. Every prediction passes the EEOC four-fifths fairness gate before it reaches the screen.'),
    ('shap',
     'Per-Employee SHAP Drivers',
     'Click any Review_ID (RV-NNNNNN) to decompose the risk score into its three top contributors. No raw Emp_ID is ever rendered.'),
    ('compare',
     'Side-by-Side Comparison',
     'Compare two pseudonymised employees across drivers, survival curves, and recommended interventions in one screen.'),
    ('causal',
     'Causal Uplift (EconML T-Learner)',
     'Per-employee CATE across three intervention arms — compensation, manager rotation, learning budget. Retention budget is no longer flat-distributed.'),
    ('rag',
     'Multilingual RAG with Provider Fallback',
     'EN/ES/HI exit-interview retrieval. 250 ms reachability probe routes to OpenAI when up, hashing fallback (256-d) when down. Render path never blocks.'),
    ('ai_assistant',
     'NL→SQL with TF-IDF Safety Net',
     'Plain-English questions answered in under 12 seconds. TF-IDF gold-question fallback over 50 vetted questions when the LLM is unreachable.'),
    ('fairness',
     'Fairness Audit (Four-Fifths Gate)',
     'Per-department disparate-impact ratios with a blocking threshold at 0.80. Recommendations are suppressed — not flagged — when the gate fails.'),
    ('alerts',
     'Slack / Teams Alert Mock',
     'High-risk Review_IDs surface to a manager channel with the top three drivers and the highest-uplift intervention pre-attached.'),
]


def add_appendix_divider(prs):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = CREAM
    bg.line.fill.background()

    # accent strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.4), prs.slide_width, Inches(0.18))
    strip.fill.solid(); strip.fill.fore_color.rgb = LAV
    strip.line.fill.background()

    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.7), Inches(1.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = 'Appendix — Dashboard Gallery'
    p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = INK

    sub = s.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.7), Inches(2.0))
    stf = sub.text_frame; stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.text = 'Eight live screenshots from the current AttriSense build (Pixel pastel theme, dark-ink axes, salted Review_IDs).'
    sp.font.size = Pt(20); sp.font.color.rgb = INK


def add_image_slide(prs, slug, headline, caption):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)

    # cream background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = CREAM
    bg.line.fill.background()

    # accent strip top
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18))
    strip.fill.solid(); strip.fill.fore_color.rgb = LAV
    strip.line.fill.background()

    # headline
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.30), Inches(12.3), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = headline
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = INK

    # image — fit inside 12.3" x 5.4" box centred horizontally below headline
    img_path = IMG / f'{slug}.png'
    if not img_path.exists():
        print(f'  ! missing {img_path}')
        return
    # Insert with width=12.0" and let height auto-scale
    pic = s.shapes.add_picture(str(img_path), Inches(0.66), Inches(1.05), width=Inches(12.0))
    # If too tall, cap height
    max_h = Inches(5.6)
    if pic.height > max_h:
        ratio = max_h / pic.height
        pic.height = max_h
        pic.width = int(pic.width * ratio)
        pic.left = int((prs.slide_width - pic.width) / 2)

    # caption at bottom
    cap = s.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.55))
    ctf = cap.text_frame; ctf.word_wrap = True
    cp = ctf.paragraphs[0]; cp.text = caption
    cp.font.size = Pt(12); cp.font.color.rgb = INK; cp.font.italic = True


def main():
    prs = Presentation(str(PPTX))
    print(f'Loaded {PPTX.name} with {len(prs.slides)} slides')
    add_appendix_divider(prs)
    for slug, head, cap in GALLERY:
        add_image_slide(prs, slug, head, cap)
        print(f'  + {slug}.png')
    prs.save(str(PPTX))
    print(f'Saved {PPTX.name} — now {len(prs.slides)} slides')


if __name__ == '__main__':
    main()
