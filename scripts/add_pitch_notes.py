"""Add speaker notes to outputs/AttriSense_VP_Pitch.pptx.

For each slide, the body bullets ARE the speaker notes (verbatim). This
script copies the body text into the slide's notes pane so the deck
opens with full talking points already populated.

Run after scripts/generate_pptx.py has produced the deck.
"""
from pathlib import Path

from pptx import Presentation

REPO = Path(__file__).resolve().parents[1]
PPTX = REPO / "outputs" / "AttriSense_VP_Pitch.pptx"

NOTES = [
    # 1
    "Open: AttriSense in one line — explainable, fair, resilient retention "
    "analytics. Audience: Engineering VP. Time: 12 minutes plus questions. "
    "Pause two seconds before clicking next.",
    # 2
    "Five sentences, deliver as a block. Land hardest on item 3 — every "
    "discipline below it is a code surface, not a slide deck. If they "
    "interrupt with 'show me', skip to slide 7 (resilience demo).",
    # 3
    "Read the $36M figure out loud — VPs anchor on dollars. Then say 'every "
    "gap below is a code surface, and we close them as system components'. "
    "Don't pause for questions here; this slide sets up the asks.",
    # 4
    "Five surfaces, seven tabs. Recite each line, then say 'all five live "
    "behind one Streamlit app — there is no second service to operate'. If "
    "they ask about T-Learner, point to slide 9 ablation.",
    # 5
    "Trace the dependency arrow top-down: app calls insights, insights call "
    "models, models call features, features call data, ops touches all "
    "five. Stress 'one entry point per layer' — that is the substitution "
    "promise. If pressed, name the entry point per layer.",
    # 6
    "Six lines, six checks. The line that earns trust with engineering VPs "
    "is 'CI runs ruff plus the full pytest suite on every push'. The line "
    "that earns trust with security is 'non-root container with healthcheck'. "
    "Mention dotenv last so you can transition into the next slide on "
    "secrets-management posture.",
    # 7
    "This is the demo slide. If a laptop is open, run the dashboard with the "
    "network cable pulled and show the green 'provider: hashing' tag in the "
    "RAG results. Otherwise, walk the bullets. End with 'ten lines of code, "
    "100% availability, not in any HR paper I have seen'.",
    # 8
    "Two ideas: model-side (four-fifths gate, audit row) and reviewer-side "
    "(Review_ID pseudonym). The Review_ID story is the senior-thinking "
    "moment — most candidates stop at 'we removed protected attributes from "
    "features'. We didn't, and we shipped a reviewer-side mitigation too.",
    # 9
    "Numbers slide. Read ROC-AUC, PR-AUC, Brier in that order; the order "
    "matters because each is more conservative than the last. Then NL-to-SQL "
    "headline: 86% with the LLM, 100% non-empty with TF-IDF fallback. "
    "Anticipated question: 'is 0.91 ROC-AUC over-fit?' — slide 10 answers it.",
    # 10
    "The ablation slide is where engineering VPs lean in. Walk row 2 — "
    "'SMOTE before split inflates ROC-AUC by 0.04, this is widespread in "
    "published HR work and we consider it a leakage bug'. That sentence is "
    "what differentiates a senior IC from a junior one.",
    # 11
    "Read each line. Then close with the verbatim phrase: 'restraint is a "
    "feature. Each item above is a PR I would close.' Pause two seconds. "
    "This is the slide that recruiters will quote back to you.",
    # 12
    "Five concrete items, each with an owner and a DoD in my head. If asked "
    "to prioritise, top two are coverage gate + intersectional fairness. "
    "Bottom one to drop is the Slack feedback loop — it is high effort and "
    "depends on a real HRIS integration we don't yet have.",
    # 13
    "Five risks, five mitigations. The one to land hardest is 'worker "
    "dignity'. A VP who has not thought about this will appreciate that we "
    "have. A VP who has thought about it will appreciate that the mitigation "
    "is encoded in the intended-use document, not just spoken.",
    # 14
    "Asks slide. Read the asks crisply, then stop talking. The first "
    "question they ask tells you which ask they will say yes to. Do not "
    "negotiate against yourself by softening any of the four asks.",
    # 15
    "Closing tagline. Read all three lines, pause, then say 'questions?' "
    "and stop. Whoever speaks first signals what they care about; let the "
    "rest of the conversation follow that signal.",
]


def main():
    prs = Presentation(str(PPTX))
    if len(prs.slides) != len(NOTES):
        print(f"WARN: {len(prs.slides)} slides but {len(NOTES)} notes")
    for slide, note in zip(prs.slides, NOTES):
        slide.notes_slide.notes_text_frame.text = note
    prs.save(str(PPTX))
    print(f"Added speaker notes to {len(NOTES)} slides in {PPTX.name}")


if __name__ == "__main__":
    main()
